"""
Document loaders for all supported file types.

Each loader reads a file and returns a Document object.
Supported types: .txt, .md, .html, .py, .pdf, .pptx

Usage:
    doc = load_document("path/to/file.pdf")
"""

from __future__ import annotations

import logging
import re
import uuid
from pathlib import Path

from src.models import Document, SourceType

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Type detection
# ---------------------------------------------------------------------------

_EXT_TO_SOURCE_TYPE: dict[str, SourceType] = {
    ".txt": SourceType.TXT,
    ".md": SourceType.MD,
    ".html": SourceType.HTML,
    ".htm": SourceType.HTML,
    ".py": SourceType.PY,
    ".pdf": SourceType.PDF,
    ".pptx": SourceType.PPTX,
}


def get_source_type(path: str | Path) -> SourceType:
    """Return the SourceType for the given file path.

    Args:
        path: Path to the file.

    Returns:
        SourceType enum value.

    Raises:
        ValueError: If the file extension is not supported.
    """
    ext = Path(path).suffix.lower()
    if ext not in _EXT_TO_SOURCE_TYPE:
        raise ValueError(f"Unsupported file extension: {ext!r}")
    return _EXT_TO_SOURCE_TYPE[ext]


# ---------------------------------------------------------------------------
# Text-based loaders
# ---------------------------------------------------------------------------

def _load_text(path: Path) -> str:
    """Read a plain text file, trying UTF-8 first then cp932 (Windows Japanese)."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        logger.debug("UTF-8 failed for %s, trying cp932", path)
        return path.read_text(encoding="cp932")


def load_txt(path: str | Path) -> Document:
    """Load a plain text file.

    Args:
        path: Path to the .txt file.

    Returns:
        Document with source_type=TXT.
    """
    p = Path(path)
    content = _load_text(p)
    logger.info("Loaded TXT: %s (%d chars)", p.name, len(content))
    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source_file=p.name,
        source_type=SourceType.TXT,
        metadata={"file_path": str(p.resolve())},
    )


def load_md(path: str | Path) -> Document:
    """Load a Markdown file.

    Args:
        path: Path to the .md file.

    Returns:
        Document with source_type=MD.
    """
    p = Path(path)
    content = _load_text(p)
    logger.info("Loaded MD: %s (%d chars)", p.name, len(content))
    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source_file=p.name,
        source_type=SourceType.MD,
        metadata={"file_path": str(p.resolve())},
    )


def load_html(path: str | Path) -> Document:
    """Load an HTML file and extract readable text using BeautifulSoup.

    Args:
        path: Path to the .html or .htm file.

    Returns:
        Document with source_type=HTML. The content is the raw HTML
        (chunkers will handle tag-based splitting).
    """
    p = Path(path)
    content = _load_text(p)
    logger.info("Loaded HTML: %s (%d chars)", p.name, len(content))
    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source_file=p.name,
        source_type=SourceType.HTML,
        metadata={"file_path": str(p.resolve())},
    )


def load_py(path: str | Path) -> Document:
    """Load a Python source file.

    Args:
        path: Path to the .py file.

    Returns:
        Document with source_type=PY.
    """
    p = Path(path)
    content = _load_text(p)
    logger.info("Loaded PY: %s (%d chars)", p.name, len(content))
    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source_file=p.name,
        source_type=SourceType.PY,
        metadata={"file_path": str(p.resolve())},
    )


# ---------------------------------------------------------------------------
# PDF loader
# ---------------------------------------------------------------------------

def load_pdf(path: str | Path, *, rendered_dir: str | Path | None = None) -> Document:
    """Load a PDF file using PyMuPDF.

    Extracts text from all pages and optionally renders each page as a PNG
    image for Vision LLM processing.

    The content field contains all page texts joined with page-break markers.
    Metadata includes per-page text and image paths (if rendered_dir given).

    Args:
        path: Path to the .pdf file.
        rendered_dir: Directory to save page images. If None, no images are
                      rendered. Required for Vision LLM ingestion.

    Returns:
        Document with source_type=PDF.
        metadata["pages"] is a list of dicts:
            {"page_number": int, "text": str, "image_path": str | None}
    """
    import fitz  # PyMuPDF

    p = Path(path)
    doc = fitz.open(str(p))

    rendered_path = Path(rendered_dir) if rendered_dir else None
    if rendered_path:
        rendered_path.mkdir(parents=True, exist_ok=True)

    pages: list[dict] = []
    full_text_parts: list[str] = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text()
        image_path: str | None = None

        if rendered_path:
            # Render page as PNG (DPI from config; default 200 here — caller
            # can pass rendered_dir only when they want images)
            matrix = fitz.Matrix(200 / 72, 200 / 72)  # 200 DPI
            pix = page.get_pixmap(matrix=matrix)
            img_file = rendered_path / f"{p.stem}_page{page_num:04d}.png"
            pix.save(str(img_file))
            image_path = str(img_file)

        pages.append({"page_number": page_num, "text": text, "image_path": image_path})
        full_text_parts.append(f"[Page {page_num}]\n{text}")

    doc.close()

    content = "\n\n".join(full_text_parts)
    logger.info("Loaded PDF: %s (%d pages, %d chars)", p.name, len(pages), len(content))

    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source_file=p.name,
        source_type=SourceType.PDF,
        metadata={
            "file_path": str(p.resolve()),
            "page_count": len(pages),
            "pages": pages,
        },
    )


# ---------------------------------------------------------------------------
# PPTX loader
# ---------------------------------------------------------------------------

def _export_slide_image_win32(pptx_path: Path, slide_index: int, output_path: Path) -> bool:
    """Export a single slide as PNG using win32com (requires PowerPoint).

    Args:
        pptx_path: Path to the .pptx file.
        slide_index: 0-based slide index.
        output_path: Path to save the PNG.

    Returns:
        True if successful, False if win32com is unavailable.
    """
    try:
        import win32com.client  # type: ignore[import]

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True  # Required on some systems
        prs = powerpoint.Presentations.Open(str(pptx_path.resolve()))
        slide = prs.Slides[slide_index + 1]  # COM is 1-based
        slide.Export(str(output_path.resolve()), "PNG")
        prs.Close()
        powerpoint.Quit()
        return True
    except Exception as exc:
        logger.debug("win32com slide export failed: %s", exc)
        return False


def load_pptx(path: str | Path, *, rendered_dir: str | Path | None = None) -> Document:
    """Load a PPTX file using python-pptx.

    Extracts text from all text frames and shapes. If rendered_dir is given
    and PowerPoint (win32com) is available, each slide is also exported as
    a PNG for Vision LLM processing.

    Args:
        path: Path to the .pptx file.
        rendered_dir: Directory to save slide images (requires PowerPoint).

    Returns:
        Document with source_type=PPTX.
        metadata["slides"] is a list of dicts:
            {"slide_number": int, "text": str, "image_path": str | None}
    """
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401 (imported for type completeness)

    p = Path(path)
    prs = Presentation(str(p))

    rendered_path = Path(rendered_dir) if rendered_dir else None
    if rendered_path:
        rendered_path.mkdir(parents=True, exist_ok=True)

    slides: list[dict] = []
    full_text_parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Extract all text from shapes (text boxes, titles, content, tables)
        shape_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        shape_texts.append(line)
            # Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        shape_texts.append(" | ".join(cells))

        text = "\n".join(shape_texts)

        image_path: str | None = None
        if rendered_path:
            img_file = rendered_path / f"{p.stem}_slide{slide_num:04d}.png"
            success = _export_slide_image_win32(p, slide_num - 1, img_file)
            if success and img_file.exists():
                image_path = str(img_file)
            else:
                logger.debug("Slide image export skipped for slide %d", slide_num)

        slides.append({"slide_number": slide_num, "text": text, "image_path": image_path})
        full_text_parts.append(f"[Slide {slide_num}]\n{text}")

    content = "\n\n".join(full_text_parts)
    logger.info("Loaded PPTX: %s (%d slides, %d chars)", p.name, len(slides), len(content))

    return Document(
        doc_id=str(uuid.uuid4()),
        content=content,
        source_file=p.name,
        source_type=SourceType.PPTX,
        metadata={
            "file_path": str(p.resolve()),
            "slide_count": len(slides),
            "slides": slides,
        },
    )


# ---------------------------------------------------------------------------
# Unified entry point
# ---------------------------------------------------------------------------

def load_document(
    path: str | Path,
    *,
    rendered_dir: str | Path | None = None,
) -> Document:
    """Load any supported document file and return a Document.

    Dispatches to the appropriate loader based on file extension.

    Args:
        path: Path to the document file.
        rendered_dir: For PDF/PPTX only — directory to save rendered page/
                      slide images for Vision LLM processing.

    Returns:
        Document ready for chunking.

    Raises:
        ValueError: If the file extension is not supported.
        FileNotFoundError: If the file does not exist.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    source_type = get_source_type(p)

    loaders = {
        SourceType.TXT: lambda: load_txt(p),
        SourceType.MD: lambda: load_md(p),
        SourceType.HTML: lambda: load_html(p),
        SourceType.PY: lambda: load_py(p),
        SourceType.PDF: lambda: load_pdf(p, rendered_dir=rendered_dir),
        SourceType.PPTX: lambda: load_pptx(p, rendered_dir=rendered_dir),
    }

    return loaders[source_type]()
